# ============================================================
# _slides/theme.py — Petit moteur de slides partagé (python-pptx)
# ============================================================
# Ce fichier est la "boîte à outils" commune à TOUS les cours d'IA.
# Chaque cours a un fichier slides/contenu.py qui décrit SES slides
# (juste des titres + des puces) et appelle ces fonctions pour
# générer un vrai fichier .pptx, avec un look identique partout.
#
# Tu n'as normalement PAS besoin de toucher ce fichier.
# Pour changer le CONTENU d'un cours -> édite son slides/contenu.py.
# Pour changer le STYLE de TOUS les cours -> édite les couleurs ci-dessous.
# ============================================================

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ------------------------------------------------------------
# 1. La charte graphique (change ça pour reskinner tous les cours)
# ------------------------------------------------------------
BLEU = RGBColor(0x1F, 0x3A, 0x5F)        # bleu nuit, pour les titres
ACCENT = RGBColor(0x2E, 0x86, 0xDE)      # bleu vif, pour les accents
GRIS = RGBColor(0x33, 0x33, 0x33)        # texte courant
GRIS_CLAIR = RGBColor(0x88, 0x88, 0x88)  # sous-titres, pieds de page
BLANC = RGBColor(0xFF, 0xFF, 0xFF)
FOND_SECTION = RGBColor(0x1F, 0x3A, 0x5F)  # fond des slides "section"

POLICE = "Calibri"

# Dimensions 16:9
LARGEUR = Inches(13.333)
HAUTEUR = Inches(7.5)


# ------------------------------------------------------------
# 2. Créer un nouveau diaporama
# ------------------------------------------------------------
def nouveau_deck():
    """Crée une présentation vide au format 16:9."""
    prs = Presentation()
    prs.slide_width = LARGEUR
    prs.slide_height = HAUTEUR
    return prs


def _slide_vide(prs):
    """Ajoute une slide totalement vierge (layout 'blank' = index 6)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _zone_texte(slide, gauche, haut, largeur, hauteur):
    box = slide.shapes.add_textbox(gauche, haut, largeur, hauteur)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


# ------------------------------------------------------------
# 3. Slide de TITRE (la première de chaque cours)
# ------------------------------------------------------------
def slide_titre(prs, titre, sous_titre=""):
    slide = _slide_vide(prs)

    # Bandeau d'accent en haut
    barre = slide.shapes.add_shape(
        1, Inches(0), Inches(2.6), LARGEUR, Inches(0.12)  # 1 = rectangle
    )
    barre.fill.solid()
    barre.fill.fore_color.rgb = ACCENT
    barre.line.fill.background()

    tf = _zone_texte(slide, Inches(1), Inches(2.9), Inches(11.3), Inches(2))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = titre
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = BLEU
    run.font.name = POLICE

    if sous_titre:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = sous_titre
        r2.font.size = Pt(22)
        r2.font.color.rgb = GRIS_CLAIR
        r2.font.name = POLICE
    return slide


# ------------------------------------------------------------
# 4. Slide de SECTION (séparateur entre les grandes parties)
# ------------------------------------------------------------
def slide_section(prs, titre):
    slide = _slide_vide(prs)

    # Fond plein bleu nuit
    fond = slide.shapes.add_shape(1, Inches(0), Inches(0), LARGEUR, HAUTEUR)
    fond.fill.solid()
    fond.fill.fore_color.rgb = FOND_SECTION
    fond.line.fill.background()
    # Renvoie le rectangle au fond
    slide.shapes._spTree.remove(fond._element)
    slide.shapes._spTree.insert(2, fond._element)

    tf = _zone_texte(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.5))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = titre
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = BLANC
    run.font.name = POLICE
    return slide


# ------------------------------------------------------------
# 5. Slide à PUCES (le cheval de bataille)
# ------------------------------------------------------------
def slide_puces(prs, titre, puces):
    """
    puces : liste d'éléments. Chaque élément est :
      - une chaîne  -> puce de niveau 1
      - un tuple (texte, niveau)  -> niveau 0, 1 ou 2 pour indenter
    """
    slide = _slide_vide(prs)

    # Titre
    tf_titre = _zone_texte(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(1))
    pt = tf_titre.paragraphs[0]
    rt = pt.add_run()
    rt.text = titre
    rt.font.size = Pt(32)
    rt.font.bold = True
    rt.font.color.rgb = BLEU
    rt.font.name = POLICE

    # Trait sous le titre
    trait = slide.shapes.add_shape(1, Inches(0.7), Inches(1.35), Inches(3), Pt(3))
    trait.fill.solid()
    trait.fill.fore_color.rgb = ACCENT
    trait.line.fill.background()

    # Corps : les puces
    tf = _zone_texte(slide, Inches(0.8), Inches(1.7), Inches(11.8), Inches(5.3))
    premier = True
    for item in puces:
        texte, niveau = (item, 0) if isinstance(item, str) else item
        p = tf.paragraphs[0] if premier else tf.add_paragraph()
        premier = False
        p.level = niveau
        p.space_after = Pt(8)
        run = p.add_run()
        # Petite puce visuelle selon le niveau
        prefixe = "• " if niveau == 0 else ("– " if niveau == 1 else "· ")
        run.text = prefixe + texte
        run.font.size = Pt(24 - 2 * niveau)
        run.font.color.rgb = GRIS if niveau == 0 else GRIS_CLAIR
        run.font.bold = niveau == 0 and texte.endswith(":")
        run.font.name = POLICE
    return slide


# ------------------------------------------------------------
# 6. Sauvegarder
# ------------------------------------------------------------
def sauver(prs, chemin):
    nb = len(prs.slides._sldIdLst)
    prs.save(chemin)
    print(f"✅ Diaporama généré : {chemin}  ({nb} slides)")
